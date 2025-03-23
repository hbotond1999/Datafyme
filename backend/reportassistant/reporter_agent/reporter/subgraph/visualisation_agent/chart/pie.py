from pptx.enum.chart import XL_LEGEND_POSITION
from pydantic import BaseModel, Field
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart
from reporter_agent.reporter.subgraph.visualisation_agent.chart.color import ColorPalette
from pptx.chart.data import ChartData

class PieChart(BaseModel, Chart):
    category_column_name: str = Field(description="This axis displays the categories or groups being compared. Each slice  corresponds to one category, such as different products, countries")
    values_column_name: str = Field(description="This axis represents the numerical values associated with each category")

    def create_meta_data(self):
        return {
            "x_axis": self.category_column_name,
            "y_axis": self.values_column_name
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        x_axis = chart.meta_data["metadata"]["x_axis"]
        y_axis = chart.meta_data["metadata"]["y_axis"]

        color_palette = ColorPalette()
        border_colors, background_colors  = color_palette.get_colors(len(data[x_axis]))
        return {
            "type": "pie",
            "data": {
                "labels": data[x_axis],
                "datasets": [{
                    "data": data[y_axis],
                    "borderWidth": 1,
                    "borderColor": border_colors,
                    "backgroundColor": background_colors,
                }]
            },
            "options": {
                "maintainAspectRatio": False,
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a pie chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """

        category_column = chart_metadata.meta_data["metadata"]["x_axis"]
        values_column = chart_metadata.meta_data["metadata"]["y_axis"]

        chart_data = ChartData()

        categories = data[category_column]
        chart_data.categories = categories

        values = data[values_column]
        chart_data.add_series('Values', values)
        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        if chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"Distribution of {values_column}"

        plot = pptx_chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = 8
        data_labels.font.bold = True

        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False

        return pptx_chart