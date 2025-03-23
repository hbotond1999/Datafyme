from pptx.chart.data import CategoryChartData
from pydantic import BaseModel, Field

from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION


class HistogramChart(BaseModel, Chart):
    y_axis_column: str = Field(description="The Y-axis displays the count or frequency of values falling within each bin. The height of each bar represents the number of data points within that bin’s range.")

    def create_meta_data(self):
        return {
            "y_axis": self.y_axis_column
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        y_axis = chart.meta_data["metadata"]["y_axis"]

        return {
            "type": "bar",
            "data": {
                "labels": [*range(len(data[y_axis]))],
                "datasets": [{
                    "label": y_axis,
                    "data": data[y_axis],
                    "backgroundColor": "rgba(75, 192, 192, 0.6)",
                    "borderColor": "rgba(75, 192, 192, 1)",
                    "borderWidth": 1
                }]
            },
            "options": {
                "maintainAspectRatio": False,
                "scales": {
                    "y": {
                        "beginAtZero": True
                    }
                }
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a histogram chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """

        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]

        chart_data = CategoryChartData()

        categories = [str(i) for i in range(len(data[y_axis]))]
        chart_data.categories = categories

        chart_data.add_series(y_axis, data[y_axis].tolist())

        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        if chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"Histogram of {y_axis}"

        category_axis = pptx_chart.category_axis
        category_axis.has_major_gridlines = False

        # Format y-axis (value axis)
        value_axis = pptx_chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.has_minor_gridlines = False
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = "Frequency"

        pptx_chart.has_legend = True
        pptx_chart.legend.position =  XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False

        plot = pptx_chart.plots[0]
        plot.gap_width = 0

        return pptx_chart