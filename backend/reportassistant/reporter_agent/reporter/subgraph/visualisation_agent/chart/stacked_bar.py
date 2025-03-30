from typing import Optional, List

from pptx.enum.chart import XL_LEGEND_POSITION
from pydantic import BaseModel, Field

from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart
from reporter_agent.reporter.subgraph.visualisation_agent.chart.color import ColorPalette
from reporter_agent.reporter.subgraph.visualisation_agent.chart.utils import axis_date_str_converter


class StackedBarChart(BaseModel, Chart):
    x_axis_column_name: str = Field( description="This x axis displays the categories or groups being compared. Each bar corresponds to one category, such as different products, countries.")
    y_value_column_name: str = Field(description="This represents the values that are being measured and visualized on the y-axis.")
    category_column_name: str = Field( description="This column defines the categories that break down the bars into segments, representing different parts of each total value.")
    date_or_date_time_format: Optional[str] = Field( description='The date or date time format to use, to show values of the x_axis_column in the chart. It should be Python strftime method compatible.', default=None)
    def create_meta_data(self):
        return {
            "x_axis": self.x_axis_column_name,
            "y_axis": self.y_value_column_name,
            "category_column_name": self.category_column_name,
            "date_format": self.date_or_date_time_format
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        date_format = chart.meta_data["metadata"].get("date_format", None)
        raw_labels = data[chart.meta_data["metadata"]["x_axis"]]
        if not date_format:
            labels =  list(set(data[chart.meta_data["metadata"]["x_axis"]]))
        else:
            labels = axis_date_str_converter(dates=list(set(data[chart.meta_data["metadata"]["x_axis"]])), date_format=date_format)

        y_axis = chart.meta_data["metadata"]["y_axis"]
        category = chart.meta_data["metadata"]["category_column_name"]

        uniques_categories = list(set(data[category]))
        color_palette = ColorPalette()
        colors_without_opacity, colors_with_opacity = color_palette.get_colors( len(uniques_categories))
        category_color_map = {item: {"backgroundColor": colors_with_opacity[i], "borderColor": colors_without_opacity[i]}  for i, item in enumerate(uniques_categories)}

        stack_data = {}
        for category_value, y_value in zip(data[category], data[y_axis]):
            if stack_data.get(category_value) is None:
                stack_data[category_value] = [y_value]
            else:
                stack_data[category_value].append(y_value)

        datasets = []
        for category_value in stack_data:
            datasets.append({
                "label": category_value,
                "data": stack_data[category_value],
                "backgroundColor": category_color_map[category_value]["backgroundColor"],
                "borderColor": category_color_map[category_value]["borderColor"],
                "borderWidth": 1,
                "stack": category_value
            })
        return {
            "type": "bar",
            "data": {
                "labels": labels,
                "datasets": datasets
            },
            "options": {
                "maintainAspectRatio": False,
                "scales": {
                    "x": {
                        "stacked": True
                    },
                    "y": {
                        "stacked": True,
                        "beginAtZero": True
                    }
                }
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a stacked bar chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.chart.data import CategoryChartData
        from pptx.dml.color import RGBColor
        import pandas as pd

        x_axis = chart_metadata.meta_data["metadata"]["x_axis"]
        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]
        category_column = chart_metadata.meta_data["metadata"]["category_column_name"]
        date_format = chart_metadata.meta_data["metadata"].get("date_format")


        pivot_df = pd.pivot_table(
            data,
            values=y_axis,
            index=x_axis,
            columns=category_column,
            aggfunc='sum'
        ).fillna(0)


        chart_data = CategoryChartData()

        if date_format:
            categories = axis_date_str_converter(dates=pivot_df.index.tolist(), date_format=date_format)
        else:
            categories = pivot_df.index.tolist()

        chart_data.categories = categories

        for i, category_value in enumerate(pivot_df.columns):
            values = pivot_df[category_value].tolist()
            chart_data.add_series(str(category_value), values)

        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        if hasattr(chart_metadata, 'title') and chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"{y_axis} by {x_axis} and {category_column}"

        category_axis = pptx_chart.category_axis
        category_axis.has_major_gridlines = False
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = x_axis

        value_axis = pptx_chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.has_minor_gridlines = False
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = y_axis

        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False
        return pptx_chart

    def validate_chart_data(self, column_names: List[str]):
        error_messages = []
        if self.x_axis_column_name not in column_names:
            error_messages.append(f"The {self.x_axis_column_name} column is not in the dataset.")

        if self.y_value_column_name not in column_names:
            error_messages.append(f"The {self.y_value_column_name} column is not in the dataset.")
        return error_messages