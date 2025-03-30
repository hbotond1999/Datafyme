from typing import List

from pptx.enum.chart import XL_LEGEND_POSITION
from pydantic import BaseModel,Field
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import XyChartData
from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart


class ScatterChart(BaseModel, Chart):
    x_axis_column_name: str = Field(description="The X-axis in a scatter chart represents the independent variable or the variable that is presumed to influence the dependent variable on the Y-axis. This is often something measurable, such as time, temperature, age, or any factor that might logically lead to changes in the variable on the Y-axis")
    y_axis_column_name: str = Field(description="he Y-axis represents the dependent variable, the variable that potentially changes in response to variations in the independent variable on the X-axis. It is often the outcome or effect being studied, like sales, test scores, or response time.")

    def create_meta_data(self):
        return {
            "x_axis": self.x_axis_column_name,
            "y_axis": self.y_axis_column_name
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        x_axis = chart.meta_data["metadata"]["x_axis"]
        y_axis = chart.meta_data["metadata"]["y_axis"]

        new_data = [{"x": x, "y": y} for x, y in zip(data[x_axis], data[y_axis])]

        return {
            "type": "scatter",
            "data": {
                "datasets": [{
                    "label": y_axis,
                    "data": new_data,
                    "backgroundColor": "rgba(255, 99, 132, 0.6)",
                    "borderColor": "rgba(255, 99, 132, 1)",
                    "borderWidth": 1
                }]
            },
            "options": {
                "maintainAspectRatio": False,
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a scatter chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """


        # Get metadata
        x_axis = chart_metadata.meta_data["metadata"]["x_axis"]
        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]

        # Create XY chart data (for scatter plots)
        chart_data = XyChartData()

        # Add a series
        series_name = chart_metadata.title if hasattr(chart_metadata, 'title') and chart_metadata.title else y_axis
        series = chart_data.add_series(series_name)

        # Add data points
        for i in range(len(data)):
            x_val = float(data[x_axis].iloc[i])
            y_val = float(data[y_axis].iloc[i])
            series.add_data_point(x_val, y_val)

        # Create the chart on the slide
        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER,
            x, y, cx, cy,
            chart_data
        ).chart

        # Set chart title
        pptx_chart.has_title = True
        if hasattr(chart_metadata, 'title') and chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"{y_axis} vs {x_axis}"

        # Format axes
        # X-axis
        category_axis = pptx_chart.category_axis
        category_axis.has_major_gridlines = True
        category_axis.has_title = True
        category_axis.axis_title.text_frame.text = x_axis

        # Y-axis
        value_axis = pptx_chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.has_title = True
        value_axis.axis_title.text_frame.text = y_axis

        # Add legend
        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False

        # Format the scatter points
        plot = pptx_chart.plots[0]
        series = plot.series[0]

        # Set marker style
        marker = series.marker
        marker.size = 10
        marker.format.fill.solid()

        return pptx_chart

    def validate_chart_data(self, column_names: List[str]):
        error_messages = []
        if self.x_axis_column_name not in column_names:
            error_messages.append(f"The {self.x_axis_column_name} column is not in the dataset.")

        if self.y_axis_column_name not in column_names:
            error_messages.append(f"The {self.y_axis_column_name} column is not in the dataset.")
        return error_messages