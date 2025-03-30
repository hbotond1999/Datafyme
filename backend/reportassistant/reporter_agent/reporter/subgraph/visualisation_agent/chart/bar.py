from typing import Optional

from pptx.chart.data import CategoryChartData
from pydantic import BaseModel, Field
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart
from reporter_agent.reporter.subgraph.visualisation_agent.chart.color import ColorPalette
from reporter_agent.reporter.subgraph.visualisation_agent.chart.utils import axis_date_str_converter


class BarChart(BaseModel, Chart):
    category_column_name: str = Field(description="This axis displays the categories or groups being compared. Each bar corresponds to one category, such as different products, countries")
    values_column_name: str = Field(description="This axis represents the numerical values associated with each category. The height (or length, in a horizontal bar chart) of each bar shows the value of that category, such as revenue, population, or any other quantitative metric.")
    date_or_date_time_format: Optional[str] = Field(description='The date or date time format to use, to show values of the x_axis_column in the chart. It should be Python strftime method compatible.', default=None)

    def create_meta_data(self):
        return {
            "x_axis": self.category_column_name,
            "y_axis": self.values_column_name,
            "date_format": self.date_or_date_time_format
        }

    @classmethod
    def create_chart_data(cls, chart, data):

        x_axis = chart.meta_data["metadata"]["x_axis"]
        y_axis = chart.meta_data["metadata"]["y_axis"]
        date_format = chart.meta_data["metadata"].get("date_format")

        if date_format:
            labels = axis_date_str_converter(dates=data[x_axis], date_format=date_format)
        else:
            labels = data[x_axis]

        color_palette = ColorPalette()
        border_colors, background_colors  = color_palette.get_colors(len(labels))
        return {
            "type": "bar",
            "data": {
                "labels": labels,
                "datasets": [{
                    "label": y_axis,
                    "data": data[y_axis],
                    "backgroundColor": background_colors,
                    "borderColor": border_colors,
                    "borderWidth": 1
                }]
            },
            "options": {
                "maintainAspectRatio": False,
                "scales": {
                    "y": {
                        "beginAtZero": True
                    }
                }
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """


        x_axis = chart_metadata.meta_data["metadata"]["x_axis"]
        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]
        date_format = chart_metadata.meta_data["metadata"].get("date_format")

        chart_data = CategoryChartData()

        if date_format:
            categories = axis_date_str_converter(dates=data[x_axis], date_format=date_format)
        else:
            categories = data[x_axis].tolist()

        chart_data.categories = categories

        values = data[y_axis].tolist()
        series_name = y_axis
        chart_data.add_series(series_name, values)

        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        if chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"{y_axis} by {x_axis}"

        value_axis = pptx_chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.has_minor_gridlines = False

        category_axis = pptx_chart.category_axis
        category_axis.has_major_gridlines = False

        return pptx_chart