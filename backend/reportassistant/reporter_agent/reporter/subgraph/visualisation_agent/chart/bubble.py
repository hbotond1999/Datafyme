from pptx.chart.data import BubbleChartData
from pydantic import BaseModel, Field
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import BubbleChartData
from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart


class BubbleChart(BaseModel, Chart):
    X_axis_column_name: str = Field(description="This column represents one variable, which determines the horizontal position of each bubble.This could be a numeric value")
    Y_axis_column_name: str = Field(description="This column represents a second variable, which determines the vertical position of each bubble. This could be another numeric value or any other metric of interest.")
    size_column_name: str = Field(description="The size of each bubble represents a third variable, indicating the value of that variable. The larger the bubble, the greater the value of this third variable. For example, it could be used to show volume, sales, population, or any other quantitative measure.")

    def create_meta_data(self):
        return {
            "x_axis": self.X_axis_column_name,
            "y_axis": self.Y_axis_column_name,
            "size": self.size_column_name
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        x_axis = chart.meta_data["metadata"]["x_axis"]
        y_axis = chart.meta_data["metadata"]["y_axis"]
        size = chart.meta_data["metadata"]["size"]

        # Create data in the format required for a bubble chart
        new_data = [{"x": x, "y": y, "r": r} for x, y, r in zip(data[x_axis], data[y_axis], data[size])]

        return {
            "type": "bubble",
            "data": {
                "datasets": [{
                    "label": chart.title,
                    "data": new_data,
                    "backgroundColor": "rgba(255, 99, 132, 0.6)",
                    "borderColor": "rgba(255, 99, 132, 1)",
                    "borderWidth": 1
                }]
            },
            "options": {
                "maintainAspectRatio": False,
                "scales": {
                    "x": {
                        "beginAtZero": True
                    },
                    "y": {
                        "beginAtZero": True
                    }
                }
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a bubble chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint bubble chart object that has been added to the slide
        """

        x_axis = chart_metadata.meta_data["metadata"]["x_axis"]
        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]
        size = chart_metadata.meta_data["metadata"]["size"]

        chart_data = BubbleChartData()

        series_name = chart_metadata.title if hasattr(chart_metadata, 'title') else "Series 1"
        series = chart_data.add_series(series_name)

        for i in range(len(data)):
            x_val = float(data[x_axis].iloc[i])
            y_val = float(data[y_axis].iloc[i])
            size_val = float(data[size].iloc[i])

            series.add_data_point(x_val, y_val, size_val)

        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BUBBLE,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        pptx_chart.chart_title.text_frame.text = f"{y_axis} vs {x_axis} (Size: {size})"

        value_axis = pptx_chart.value_axis
        value_axis.has_major_gridlines = True
        value_axis.has_minor_gridlines = False

        category_axis = pptx_chart.category_axis
        category_axis.has_major_gridlines = True

        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False

        return pptx_chart