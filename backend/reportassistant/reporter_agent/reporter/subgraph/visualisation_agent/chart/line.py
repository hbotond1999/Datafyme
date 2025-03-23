from typing import Optional

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pydantic import BaseModel, Field
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE
from reporter_agent.reporter.subgraph.visualisation_agent.chart.abc import Chart

from reporter_agent.reporter.subgraph.visualisation_agent.chart.utils import axis_date_str_converter


class LineChart(BaseModel, Chart):
    x_axis_column_name: str = Field(description='Represents the independent variable, often time (e.g., days, months, years).')
    y_axis_column_name: str = Field(description='Shows the dependent variable, the value being measured (e.g., temperature, revenue, or population).')
    date_format: Optional[str] = Field(description='The date or date time format to use, to show values of the x_axis_column in the chart. It should be Python strftime method compatible.', default=None)

    def create_meta_data(self):
        return {
            "x_axis": self.x_axis_column_name,
            "y_axis": self.y_axis_column_name,
            "date_format": self.date_format,
        }

    @classmethod
    def create_chart_data(cls, chart, data):
        x_axis = chart.meta_data["metadata"]["x_axis"]
        y_axis = chart.meta_data["metadata"]["y_axis"]
        date_format = chart.meta_data["metadata"].get("date_format", None)
        print(date_format)
        if not date_format:
            labels = data[x_axis]
        else:
            labels = axis_date_str_converter(dates=data[x_axis], date_format=date_format)

        return {
            "type": "line",
            "data": {
                "labels": labels,
                "datasets": [{
                    "label": y_axis,
                    "data": data[y_axis],
                    "borderColor": "rgba(75, 192, 192, 1)",
                    "backgroundColor": "rgba(75, 192, 192, 0.4)",
                    "fill": True
                }]
            },
            "options": {
                "maintainAspectRatio": False,
                "scales": {
                    "x": {
                        "beginAtZero": True
                    },
                    "y": {
                        "beginAtZero": True
                    }
                }
            }
        }

    @classmethod
    def create_pptx_chart(cls, chart_metadata, data, slide, x, y, cx, cy):
        """
        Creates a line chart object that can be directly added to a PowerPoint slide.

        Args:
            chart_metadata: The chart object containing metadata
            data: The DataFrame containing the data to be plotted
            slide: The PowerPoint slide to add the chart to
            x: The x-coordinate for the chart position
            y: The y-coordinate for the chart position
            cx: The width of the chart
            cy: The height of the chart

        Returns:
            Chart: A PowerPoint chart object that has been added to the slide
        """

        x_axis = chart_metadata.meta_data["metadata"]["x_axis"]
        y_axis = chart_metadata.meta_data["metadata"]["y_axis"]
        date_format = chart_metadata.meta_data["metadata"].get("date_format")

        chart_data = CategoryChartData()

        if not date_format:
            categories = data[x_axis].tolist()
        else:
            categories = axis_date_str_converter(dates=data[x_axis], date_format=date_format)

        chart_data.categories = categories
        chart_data.add_series(y_axis, data[y_axis].tolist())

        pptx_chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            x, y, cx, cy,
            chart_data
        ).chart

        pptx_chart.has_title = True
        if hasattr(chart_metadata, 'title') and chart_metadata.title:
            pptx_chart.chart_title.text_frame.text = chart_metadata.title
        else:
            pptx_chart.chart_title.text_frame.text = f"{y_axis} over {x_axis}"

        pptx_chart.has_legend = True
        pptx_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        pptx_chart.legend.include_in_layout = False

        plot = pptx_chart.plots[0]
        series = plot.series[0]

        series.format.line.color.rgb = RGBColor(65, 105, 225)
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = 5

        return pptx_chart